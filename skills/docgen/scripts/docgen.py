#!/usr/bin/env python3
"""
Document Generation via AWS Bedrock AgentCore Runtime.

Generates pptx, docx, pdf, xlsx, or html files using a remote
AgentCore Runtime agent (Claude + Code Interpreter).

Usage:
    python3 docgen.py --type pptx --prompt "Create a 10-slide deck about AI" --output slides.pptx
    python3 docgen.py --type docx --prompt "Write a technical report" --output report.docx
    python3 docgen.py --type pptx --prompt-file prompt.txt --output deck.pptx

Environment variables:
    AWS_REGION      - AWS region (default: ap-northeast-1)
    RUNTIME_ARN     - AgentCore Runtime ARN (required)
    AWS credentials - From ~/.aws/credentials, env vars, or IAM role

Optionally reads .env from ~/codes/document-generation-mcp/.env if present.
"""

import argparse
import base64
import json
import os
import sys
import time
from pathlib import Path

VALID_TYPES = {"docx", "pdf", "pptx", "xlsx", "frontend-design"}

# Transient error patterns that warrant a retry
_TRANSIENT_PATTERNS = ("timeout", "timed out", "connection", "prematurely", "reset by peer")

# Try loading .env from document-generation-mcp project (optional)
for _env_candidate in [
    Path.home() / "codes" / "document-generation-mcp" / ".env",
    Path(__file__).resolve().parent.parent / ".env",
]:
    if _env_candidate.exists():
        for line in _env_candidate.read_text().splitlines():
            line = line.strip()
            if line and not line.startswith("#") and "=" in line:
                key, _, value = line.partition("=")
                os.environ.setdefault(key.strip(), value.strip())
        break

DEFAULT_REGION = os.environ.get("AWS_REGION", "ap-northeast-1")
DEFAULT_RUNTIME_ARN = os.environ.get("RUNTIME_ARN", "")


def _is_transient_error(error: Exception) -> bool:
    """Check if an error is transient and worth retrying."""
    err_str = str(error).lower()
    return any(p in err_str for p in _TRANSIENT_PATTERNS)


def _estimate_timeout(skill_type: str, prompt: str) -> int:
    """Estimate a reasonable timeout based on document type and complexity.

    Returns timeout in seconds.
    """
    # Base timeouts by type
    base = {"frontend-design": 120, "xlsx": 300, "docx": 300, "pdf": 300, "pptx": 600}
    timeout = base.get(skill_type, 600)

    if skill_type == "pptx":
        # Estimate slide count from prompt
        slide_count = _count_slides_in_prompt(prompt)
        if slide_count > 15:
            timeout = 1800  # 30 min for complex decks
        elif slide_count > 10:
            timeout = 1200  # 20 min
        elif slide_count > 5:
            timeout = 900   # 15 min

    # Scale up for very long prompts (more detail = more work)
    if len(prompt) > 3000:
        timeout = int(timeout * 1.3)

    return timeout


def _count_slides_in_prompt(prompt: str) -> int:
    """Estimate the number of slides requested in a prompt."""
    import re

    # Look for explicit slide count: "15-slide", "15 slides", "Create a 15-slide"
    match = re.search(r"(\d+)[\s-]*slide", prompt, re.IGNORECASE)
    if match:
        return int(match.group(1))

    # Count "Slide N:" or "第N张:" patterns — find highest number
    slide_refs = re.findall(r"(?:slide|幻灯片|第)\s*(\d+)", prompt, re.IGNORECASE)
    if slide_refs:
        return max(int(s) for s in slide_refs)

    # Count lines that look like slide definitions
    slide_lines = re.findall(r"^(?:Slide|第)\s*\d+", prompt, re.MULTILINE | re.IGNORECASE)
    if slide_lines:
        return len(slide_lines)

    return 10  # default assumption


def _split_pptx_prompt(prompt: str, max_slides_per_batch: int = 11):
    """Split a PPTX prompt into batches of ≤max_slides_per_batch slides.

    Returns a list of (batch_prompt, slide_range) tuples.
    If the prompt has ≤max_slides_per_batch slides, returns [(prompt, None)].
    """
    import re

    slide_count = _count_slides_in_prompt(prompt)
    if slide_count <= max_slides_per_batch:
        return [(prompt, None)]

    # Find slide definition lines: "Slide N:" or "第N张:" patterns
    slide_pattern = re.compile(
        r"^((?:Slide|幻灯片|第)\s*\d+\s*[:：])",
        re.MULTILINE | re.IGNORECASE
    )
    matches = list(slide_pattern.finditer(prompt))

    if len(matches) < 2:
        # Can't reliably split — return as-is
        return [(prompt, None)]

    # Split into batches
    batches = []
    # Extract header (everything before first slide definition)
    header = prompt[:matches[0].start()].rstrip()

    for batch_start in range(0, len(matches), max_slides_per_batch):
        batch_end = min(batch_start + max_slides_per_batch, len(matches))
        start_pos = matches[batch_start].start()
        end_pos = matches[batch_end].start() if batch_end < len(matches) else len(prompt)

        slide_content = prompt[start_pos:end_pos].rstrip()
        # Renumber slides in the batch starting from 1
        batch_slide_count = batch_end - batch_start
        batch_prompt = f"{header}\n\nCreate a {batch_slide_count}-slide presentation.\n\n{slide_content}"
        slide_range = f"{batch_start + 1}-{batch_end}"
        batches.append((batch_prompt, slide_range))

    return batches


def _merge_pptx_files(pptx_paths: list, output_path: str) -> bool:
    """Merge multiple PPTX files into one using python-pptx.

    Returns True on success.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu
        import copy

        if not pptx_paths:
            return False

        if len(pptx_paths) == 1:
            import shutil
            shutil.copy2(pptx_paths[0], output_path)
            return True

        # Use first file as base
        prs = Presentation(pptx_paths[0])

        for path in pptx_paths[1:]:
            src = Presentation(path)
            for slide in src.slides:
                # Add a blank slide and copy content
                layout = prs.slide_layouts[6]  # blank layout
                new_slide = prs.slides.add_slide(layout)

                # Copy shapes
                for shape in slide.shapes:
                    el = copy.deepcopy(shape._element)
                    new_slide.shapes._spTree.append(el)

                # Copy notes
                if slide.has_notes_slide:
                    notes_slide = new_slide.notes_slide
                    notes_slide.notes_text_frame.text = slide.notes_slide.notes_text_frame.text

        prs.save(output_path)
        return True

    except ImportError:
        print("  Warning: python-pptx not installed, cannot merge. Keeping individual files.")
        return False
    except Exception as e:
        print(f"  Warning: merge failed ({e}). Keeping individual files.")
        return False


def generate_pptx_split(prompt, output_path, region=None, runtime_arn=None,
                        read_timeout=None, max_retries=3, max_slides_per_batch=11):
    """Generate a large PPTX by splitting into batches and merging.

    Automatically splits prompts with >max_slides_per_batch slides into parallel
    generation jobs, then merges the results.

    Returns same dict format as generate().
    """
    from concurrent.futures import ThreadPoolExecutor, as_completed

    batches = _split_pptx_prompt(prompt, max_slides_per_batch)

    if len(batches) == 1:
        # No split needed
        return generate("pptx", prompt, output_path, region, runtime_arn,
                        read_timeout, max_retries)

    print(f"  Auto-splitting into {len(batches)} batches ({max_slides_per_batch} slides each)")

    # Generate parts in parallel
    base_name = os.path.splitext(output_path)[0]
    part_paths = []
    results = {}

    start = time.time()

    def gen_batch(idx, batch_prompt, slide_range):
        part_path = f"{base_name}_Part{idx + 1}.pptx"
        result = generate("pptx", batch_prompt, part_path, region, runtime_arn,
                          read_timeout, max_retries)
        return idx, part_path, result

    with ThreadPoolExecutor(max_workers=len(batches)) as executor:
        futures = []
        for idx, (batch_prompt, slide_range) in enumerate(batches):
            print(f"  Batch {idx + 1}: slides {slide_range}")
            futures.append(executor.submit(gen_batch, idx, batch_prompt, slide_range))

        for future in as_completed(futures):
            idx, part_path, result = future.result()
            results[idx] = (part_path, result)

    # Check results
    all_success = all(r[1]["success"] for r in results.values())
    elapsed = time.time() - start

    if not all_success:
        failed = [f"Batch {i+1}: {r[1].get('error', 'unknown')}"
                  for i, r in results.items() if not r[1]["success"]]
        return {"success": False, "error": f"Some batches failed: {'; '.join(failed)}", "elapsed": elapsed}

    # Merge parts
    part_paths = [results[i][0] for i in sorted(results.keys())]
    print(f"  All batches complete. Merging {len(part_paths)} files...")

    if _merge_pptx_files(part_paths, output_path):
        total_size = os.path.getsize(output_path)
        return {
            "success": True,
            "path": output_path,
            "size": total_size,
            "elapsed": elapsed,
            "parts": part_paths,
        }
    else:
        # Merge failed — return the individual parts as a partial success
        return {
            "success": True,
            "path": part_paths[0],
            "size": sum(os.path.getsize(p) for p in part_paths),
            "elapsed": elapsed,
            "parts": part_paths,
            "merge_failed": True,
        }


def generate(skill_type, prompt, output_path, region=None, runtime_arn=None,
             read_timeout=None, max_retries=3, persist_s3=False):
    """Generate a document via AgentCore Runtime.

    Args:
        skill_type: One of docx, pdf, pptx, xlsx, frontend-design
        prompt: Description of the document to create
        output_path: Where to save the generated file
        region: AWS region (default from env)
        runtime_arn: AgentCore Runtime ARN (default from env)
        read_timeout: HTTP read timeout in seconds (auto-estimated if None)
        max_retries: Number of retries for transient errors (default 3)
        persist_s3: Also persist output to S3 cls-laptop/slides/ (default False)

    Returns:
        dict with keys: success (bool), path (str), size (int), elapsed (float)
    """
    import boto3
    from botocore.config import Config

    region = region or DEFAULT_REGION
    runtime_arn = runtime_arn or DEFAULT_RUNTIME_ARN

    if not runtime_arn:
        return {"success": False, "error": "RUNTIME_ARN not configured. Set via env var or --runtime-arn."}

    if skill_type not in VALID_TYPES:
        return {"success": False, "error": f"Invalid type: {skill_type}. Valid: {sorted(VALID_TYPES)}"}

    # Auto-estimate timeout if not provided
    if read_timeout is None:
        read_timeout = _estimate_timeout(skill_type, prompt)

    client = boto3.client(
        "bedrock-agentcore",
        region_name=region,
        config=Config(read_timeout=read_timeout, connect_timeout=30),
    )

    payload_dict = {
        "skill_type": skill_type,
        "prompt": prompt,
        "filename": os.path.basename(output_path),
    }
    if persist_s3:
        payload_dict["persist_s3"] = True

    payload = json.dumps(payload_dict)

    start = time.time()
    last_error = None

    for attempt in range(max_retries):
        try:
            if attempt > 0:
                backoff = 10 * (attempt + 1)  # 20s, 30s
                print(f"  Retry {attempt}/{max_retries - 1} after {backoff}s backoff...")
                time.sleep(backoff)

            response = client.invoke_agent_runtime(
                agentRuntimeArn=runtime_arn,
                payload=payload,
                contentType="application/json",
                accept="application/json",
            )

            body = response.get("response")
            if body is None:
                return {"success": False, "error": "Empty response from AgentCore", "elapsed": time.time() - start}

            data = body.read() if hasattr(body, "read") else body
            if isinstance(data, bytes):
                data = data.decode("utf-8")

            result = json.loads(data)
            elapsed = time.time() - start

            if result.get("status") == "success" and "file_base64" in result:
                file_bytes = base64.b64decode(result["file_base64"])
                os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
                with open(output_path, "wb") as f:
                    f.write(file_bytes)
                ret = {
                    "success": True,
                    "path": output_path,
                    "size": len(file_bytes),
                    "elapsed": elapsed,
                }
                if result.get("s3_uri"):
                    ret["s3_uri"] = result["s3_uri"]
                return ret
            else:
                return {
                    "success": False,
                    "error": result.get("error", "No file in response"),
                    "elapsed": elapsed,
                }

        except Exception as e:
            last_error = e
            if _is_transient_error(e) and attempt < max_retries - 1:
                print(f"  Transient error: {e}")
                continue
            return {"success": False, "error": str(e), "elapsed": time.time() - start}

    return {"success": False, "error": f"All {max_retries} attempts failed: {last_error}", "elapsed": time.time() - start}


def main():
    parser = argparse.ArgumentParser(description="Generate documents via AgentCore Runtime")
    parser.add_argument("--type", dest="skill_type", choices=sorted(VALID_TYPES), required=True,
                        help="Document type to generate")
    parser.add_argument("--prompt", help="Document description")
    parser.add_argument("--prompt-file", help="Read prompt from file")
    parser.add_argument("--output", "-o", required=True, help="Output file path")
    parser.add_argument("--region", default=DEFAULT_REGION, help=f"AWS region (default: {DEFAULT_REGION})")
    parser.add_argument("--runtime-arn", default=DEFAULT_RUNTIME_ARN, help="AgentCore Runtime ARN")
    parser.add_argument("--timeout", type=int, default=None, help="Read timeout in seconds (auto-estimated if omitted)")
    parser.add_argument("--no-split", action="store_true", help="Disable auto-split for large PPTX")
    parser.add_argument("--max-slides-per-batch", type=int, default=11, help="Max slides per batch when splitting (default: 11)")
    parser.add_argument("--persist-s3", action="store_true", help="Also persist output to S3 cls-laptop/slides/")

    args = parser.parse_args()

    if args.prompt_file:
        prompt = Path(args.prompt_file).read_text().strip()
    elif args.prompt:
        prompt = args.prompt
    else:
        print("Error: --prompt or --prompt-file required", file=sys.stderr)
        sys.exit(1)

    print(f"Generating {args.skill_type} \u2192 {args.output}")
    print(f"Prompt length: {len(prompt)} chars")

    if args.persist_s3:
        print("S3 persist enabled: output will also be saved to s3://cls-laptop/slides/")

    # Auto-split for large PPTX unless disabled
    if args.skill_type == "pptx" and not args.no_split:
        slide_count = _count_slides_in_prompt(prompt)
        if slide_count > args.max_slides_per_batch:
            print(f"Detected {slide_count} slides (>{args.max_slides_per_batch}), using auto-split mode")
            result = generate_pptx_split(
                prompt=prompt,
                output_path=args.output,
                region=args.region,
                runtime_arn=args.runtime_arn,
                read_timeout=args.timeout,
                max_retries=3,
                max_slides_per_batch=args.max_slides_per_batch,
            )
        else:
            timeout = args.timeout or _estimate_timeout("pptx", prompt)
            print(f"Estimated timeout: {timeout}s")
            result = generate(
                skill_type="pptx",
                prompt=prompt,
                output_path=args.output,
                region=args.region,
                runtime_arn=args.runtime_arn,
                read_timeout=args.timeout,
                persist_s3=args.persist_s3,
            )
    else:
        if not args.timeout:
            timeout = _estimate_timeout(args.skill_type, prompt)
            print(f"Estimated timeout: {timeout}s")
        result = generate(
            skill_type=args.skill_type,
            prompt=prompt,
            output_path=args.output,
            region=args.region,
            runtime_arn=args.runtime_arn,
            read_timeout=args.timeout,
            persist_s3=args.persist_s3,
        )

    if result["success"]:
        size_kb = result["size"] / 1024
        mins, secs = divmod(int(result["elapsed"]), 60)
        print(f"Saved: {result['path']} ({size_kb:.1f} KB) in {mins}m {secs}s")
        if result.get("s3_uri"):
            print(f"S3: {result['s3_uri']}")
        if result.get("parts"):
            print(f"Parts: {', '.join(result['parts'])}")
        if result.get("merge_failed"):
            print("  Note: merge failed, individual part files preserved")
    else:
        print(f"Failed: {result['error']}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
